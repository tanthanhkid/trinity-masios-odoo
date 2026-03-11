"""Generate testcase-slides.pptx for Masi OS OpenClaw AI Agent testing."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Constants
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ORANGE = RGBColor(0xE8, 0x61, 0x3C)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BG2 = RGBColor(0x16, 0x21, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_TEXT = RGBColor(0x99, 0x99, 0x99)
GREEN = RGBColor(0x27, 0xAE, 0x60)
FONT = "Calibri"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=12,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size=12,
                          color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                          font_name=FONT, line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, text_lines,
                     font_size=12, font_color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                     font_name=FONT, line_spacing=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Adjust corner rounding
    shape.adjustments[0] = 0.05

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)

    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.alignment = alignment
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return shape


def add_avatar(slide, left, top, emoji, is_bot=True):
    size = Inches(0.4)
    bg_color = LIGHT_GRAY if is_bot else ORANGE
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = emoji
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_test_header(slide, number, title):
    # Orange circle with number
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.5), Inches(0.3), Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = ORANGE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title text
    add_textbox(slide, Inches(1.15), Inches(0.3), Inches(10), Inches(0.5),
                title, font_size=24, color=DARK_TEXT, bold=True)

    # Thin orange line under header
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()


def add_chat_pair(slide, user_text, bot_lines, user_y=1.2, bot_y=None,
                  user_font=13, bot_font=11, bot_height=None):
    sw = Inches(13.333)

    # User bubble (right side)
    user_lines = user_text if isinstance(user_text, list) else [user_text]
    user_w = Inches(6.5)
    user_h = Inches(0.2 + 0.25 * len(user_lines))
    user_left = sw - user_w - Inches(0.9)

    add_avatar(slide, sw - Inches(0.8), Inches(user_y), "\U0001F464", is_bot=False)
    add_rounded_rect(slide, user_left, Inches(user_y), user_w, user_h,
                     ORANGE, user_lines, font_size=user_font, font_color=WHITE,
                     alignment=PP_ALIGN.RIGHT)

    # Bot bubble (left side)
    if bot_y is None:
        bot_y = user_y + user_h / Inches(1) + 0.2
    if bot_height is None:
        bot_height = Inches(0.3 + 0.2 * len(bot_lines))
    bot_w = Inches(8.5)

    add_avatar(slide, Inches(0.4), Inches(bot_y), "\U0001F99E", is_bot=True)
    add_rounded_rect(slide, Inches(0.9), Inches(bot_y), bot_w, bot_height,
                     LIGHT_GRAY, bot_lines, font_size=bot_font, font_color=DARK_TEXT,
                     alignment=PP_ALIGN.LEFT, line_spacing=15)


def make_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    # Lobster emoji
    add_textbox(slide, Inches(5.5), Inches(0.4), Inches(2.5), Inches(1),
                "\U0001F99E", font_size=60, alignment=PP_ALIGN.CENTER)

    # MASI OS
    add_textbox(slide, Inches(3), Inches(1.5), Inches(7.333), Inches(0.6),
                "MASI OS", font_size=28, color=ORANGE, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Title
    add_textbox(slide, Inches(1.5), Inches(2.3), Inches(10.333), Inches(1),
                "Ki\u1ec3m th\u1eed AI Agent Odoo \u2014 OpenClaw",
                font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(2), Inches(3.3), Inches(9.333), Inches(0.6),
                "Masi OS Business Agent \u2014 K\u1ebft n\u1ed1i Odoo 18 qua MCP",
                font_size=18, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

    # Badges
    badge_texts = ["\U0001F9EA 10 Test Cases", "\U0001F527 24 MCP Tools", "\U0001F916 2 Telegram Bots"]
    badge_w = Inches(2.8)
    gap = Inches(0.3)
    total_w = 3 * badge_w + 2 * gap
    start_x = (SLIDE_W - total_w) / 2

    for i, txt in enumerate(badge_texts):
        x = start_x + i * (badge_w + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.0), badge_w, Inches(0.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4A)
        shape.line.color.rgb = ORANGE
        shape.line.width = Pt(1.5)
        shape.adjustments[0] = 0.15
        tf = shape.text_frame
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        p = tf.paragraphs[0]
        p.text = txt
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.name = FONT
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def make_test_slide(prs, number, title, user_text, bot_lines, bot_height=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_test_header(slide, number, title)
    bh = bot_height or Inches(0.3 + 0.2 * len(bot_lines))
    add_chat_pair(slide, user_text, bot_lines, bot_height=bh)
    return slide


def make_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Title
    add_textbox(slide, Inches(1), Inches(0.25), Inches(11.333), Inches(0.7),
                "\U0001F3C1 T\u1ed5ng k\u1ebft ki\u1ec3m th\u1eed",
                font_size=30, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    # Results table
    tests = [
        ("TC1", "K\u1ebft n\u1ed1i h\u1ec7 th\u1ed1ng"),
        ("TC2", "Dashboard KPIs"),
        ("TC3", "CRM Pipeline"),
        ("TC4", "T\u1ea1o kh\xe1ch h\xe0ng"),
        ("TC5", "Ki\u1ec3m tra c\xf4ng n\u1ee3"),
        ("TC6", "\u0110\u1ed5i ph\xe2n lo\u1ea1i + credit"),
        ("TC7", "Li\u1ec7t k\xea s\u1ea3n ph\u1ea9m"),
        ("TC8", "T\u1ed5ng quan kinh doanh"),
        ("TC9", "T\xecm kh\xe1ch h\xe0ng"),
        ("TC10", "V\u01b0\u1ee3t h\u1ea1n m\u1ee9c"),
    ]

    col1_x = Inches(2.5)
    col2_x = Inches(7.5)
    start_y = Inches(1.1)
    row_h = Inches(0.38)

    for i, (tc, desc) in enumerate(tests):
        y = start_y + i * row_h
        # TC label + description
        add_textbox(slide, col1_x, y, Inches(4.5), row_h,
                    f"{tc} \u2014 {desc}", font_size=13, color=WHITE)
        # PASS badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, col2_x, y + Inches(0.04),
            Inches(0.9), Inches(0.3)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN
        badge.line.fill.background()
        badge.adjustments[0] = 0.25
        tf = badge.text_frame
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        p.text = "PASS"
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Result text
    result_y = start_y + 10 * row_h + Inches(0.15)
    add_textbox(slide, Inches(2.5), result_y, Inches(8), Inches(0.5),
                "K\u1ebft qu\u1ea3: 10/10 PASS", font_size=22, color=GREEN,
                bold=True, alignment=PP_ALIGN.CENTER)

    # MCP Tools tested
    tools_y = result_y + Inches(0.6)
    tools_text = ("MCP Tools tested: odoo_server_info, odoo_dashboard_kpis, "
                  "odoo_pipeline_by_stage, odoo_crm_lead_summary, odoo_create_customer, "
                  "odoo_customer_credit_status, odoo_customer_set_classification, "
                  "odoo_write, odoo_search_read, odoo_customers_exceeding_credit")
    add_textbox(slide, Inches(1), tools_y, Inches(11.333), Inches(0.6),
                tools_text, font_size=9, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

    # CTA box
    cta_y = tools_y + Inches(0.65)
    cta_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), cta_y, Inches(8.333), Inches(1.2)
    )
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x4A)
    cta_box.line.color.rgb = ORANGE
    cta_box.line.width = Pt(1.5)
    cta_box.adjustments[0] = 0.08

    tf = cta_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "\U0001F99E H\xe3y chat th\u1eed v\u1edbi bot tr\xean Telegram!"
    p.font.size = Pt(16)
    p.font.color.rgb = ORANGE
    p.font.bold = True
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "@hdxthanhtt4bot \u2014 Bot 1 (CEO Dashboard)"
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE
    p2.font.name = FONT
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "@MASIBIO_bot \u2014 Bot 2 (Sales Agent)"
    p3.font.size = Pt(12)
    p3.font.color.rgb = WHITE
    p3.font.name = FONT
    p3.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 0: Cover
    make_cover(prs)

    # Slide 1: TC1
    make_test_slide(prs, 1, "Ki\u1ec3m tra k\u1ebft n\u1ed1i h\u1ec7 th\u1ed1ng",
        "Ki\u1ec3m tra k\u1ebft n\u1ed1i h\u1ec7 th\u1ed1ng Odoo, cho t\xf4i th\xf4ng tin server",
        [
            "K\u1ebft n\u1ed1i h\u1ec7 th\u1ed1ng Odoo \u0111ang ho\u1ea1t \u0111\u1ed9ng t\u1ed1t! \u2705",
            "",
            "Th\xf4ng tin server:",
            "\u2022 Tr\u1ea1ng th\xe1i: \u0110\xe3 k\u1ebft n\u1ed1i",
            "\u2022 Phi\xean b\u1ea3n Odoo: 18.0",
            "",
            "H\u1ec7 th\u1ed1ng \u0111\xe3 s\u1eb5n s\xe0ng \u0111\u1ec3 s\u1eed d\u1ee5ng.",
        ], bot_height=Inches(2.2))

    # Slide 2: TC2
    make_test_slide(prs, 2, "Dashboard KPIs cho CEO",
        "Cho t\xf4i xem KPI t\u1ed5ng h\u1ee3p: doanh thu, c\xf4ng n\u1ee3, pipeline, leads m\u1edbi",
        [
            "\U0001F4CA Dashboard KPIs - Th\xe1ng 3/2026",
            "",
            "\U0001F4B0 Doanh thu th\xe1ng:      0 VND",
            "\U0001F4C8 Pipeline Value:          0 VND",
            "\U0001F4B3 T\u1ed5ng c\xf4ng n\u1ee3:           0 VND",
            "\U0001F195 Leads m\u1edbi:                3 leads",
            "",
            "\U0001F504 Pipeline theo giai \u0111o\u1ea1n:",
            "  New: 3 | Qualified: 0 | Proposition: 0 | Won: 0",
            "",
            "\u26A0\uFE0F Nh\u1eadn x\xe9t: H\u1ec7 th\u1ed1ng \u0111ang \u1edf giai \u0111o\u1ea1n kh\u1edfi t\u1ea1o.",
            "3 leads \u1edf giai \u0111o\u1ea1n New c\u1ea7n qualify.",
        ], bot_height=Inches(3.5))

    # Slide 3: TC3
    make_test_slide(prs, 3, "CRM Pipeline",
        "Li\u1ec7t k\xea pipeline CRM theo t\u1eebng giai \u0111o\u1ea1n, c\xf3 bao nhi\xeau deals?",
        [
            "\U0001F504 Pipeline CRM - T\u1ed5ng h\u1ee3p theo giai \u0111o\u1ea1n",
            "",
            "\U0001F195 New:              3 deals  |  0 VND",
            "\u2705 Qualified:        0 deals  |  0 VND",
            "\U0001F4CB Proposition:    0 deals  |  0 VND",
            "\U0001F3C6 Won:              0 deals  |  0 VND",
            "",
            "\U0001F4CA T\u1ed5ng k\u1ebft:",
            "\u2022 T\u1ed5ng s\u1ed1 deals: 3",
            "\u2022 T\u1ed5ng gi\xe1 tr\u1ecb pipeline: 0 VND",
            "",
            "\u0110\u1ec1 xu\u1ea5t: 3 deals \u1edf New c\u1ea7n \u0111\u01b0\u1ee3c \u0111\xe1nh gi\xe1",
            "v\xe0 chuy\u1ec3n sang Qualified.",
        ], bot_height=Inches(3.5))

    # Slide 4: TC4
    make_test_slide(prs, 4, "T\u1ea1o kh\xe1ch h\xe0ng m\u1edbi",
        "T\u1ea1o kh\xe1ch h\xe0ng m\u1edbi: C\xf4ng ty ABC Tech, S\u0110T 0901234567, email contact@abctech.vn",
        [
            "\u2705 \u0110\xe3 t\u1ea1o kh\xe1ch h\xe0ng m\u1edbi th\xe0nh c\xf4ng!",
            "",
            "\U0001F194 ID:              14",
            "\U0001F4DB T\xean:            C\xf4ng ty ABC Tech",
            "\U0001F4DE S\u0110T:           0901234567",
            "\U0001F4E7 Email:         contact@abctech.vn",
            "\U0001F3E2 Lo\u1ea1i:           Company",
            "\U0001F3F7\uFE0F Ph\xe2n lo\u1ea1i:   new (kh\xe1ch h\xe0ng m\u1edbi)",
            "",
            "\u26A0\uFE0F KH m\u1edbi ch\u01b0a \u0111\u01b0\u1ee3c c\u1ea5p h\u1ea1n m\u1ee9c c\xf4ng n\u1ee3.",
            "\u0110\u1ec3 c\u1ea5p, c\u1ea7n \u0111\u1ed5i ph\xe2n lo\u1ea1i sang old.",
        ], bot_height=Inches(3.3))

    # Slide 5: TC5
    make_test_slide(prs, 5, "Ki\u1ec3m tra c\xf4ng n\u1ee3 kh\xe1ch h\xe0ng",
        "Ki\u1ec3m tra t\xecnh tr\u1ea1ng c\xf4ng n\u1ee3 c\u1ee7a kh\xe1ch h\xe0ng ABC Tech",
        [
            "\U0001F4B3 C\xf4ng n\u1ee3 - C\xf4ng ty ABC Tech",
            "",
            "\U0001F3F7\uFE0F Ph\xe2n lo\u1ea1i:                     new (kh\xe1ch h\xe0ng m\u1edbi)",
            "\u2705 \u0110\u01b0\u1ee3c ph\xe9p c\xf4ng n\u1ee3:        \u274C Kh\xf4ng",
            "\U0001F4B0 H\u1ea1n m\u1ee9c t\xedn d\u1ee5ng:          0 VND",
            "\U0001F4CA N\u1ee3 hi\u1ec7n t\u1ea1i:                     0 VND",
            "\U0001F4B5 T\xedn d\u1ee5ng kh\u1ea3 d\u1ee5ng:         0 VND",
            "\u26A0\uFE0F V\u01b0\u1ee3t h\u1ea1n m\u1ee9c:                Kh\xf4ng",
            "",
            "\U0001F4DD KH m\u1edbi ch\u01b0a \u0111\u01b0\u1ee3c ph\xe9p c\xf4ng n\u1ee3.",
        ], bot_height=Inches(3.1))

    # Slide 6: TC6
    make_test_slide(prs, 6, "\u0110\u1ed5i ph\xe2n lo\u1ea1i + set credit limit",
        "\u0110\u1ed5i kh\xe1ch h\xe0ng ABC Tech sang ph\xe2n lo\u1ea1i old, r\u1ed3i set h\u1ea1n m\u1ee9c c\xf4ng n\u1ee3 50 tri\u1ec7u",
        [
            "\u2705 \u0110\xe3 c\u1eadp nh\u1eadt th\xe0nh c\xf4ng!",
            "",
            "\U0001F4B3 C\xf4ng n\u1ee3 sau c\u1eadp nh\u1eadt:",
            "\U0001F3F7\uFE0F Ph\xe2n lo\u1ea1i:                     old \u2713",
            "\u2705 \u0110\u01b0\u1ee3c ph\xe9p c\xf4ng n\u1ee3:        \u2705 C\xf3",
            "\U0001F4B0 H\u1ea1n m\u1ee9c t\xedn d\u1ee5ng:          50,000,000 VND",
            "\U0001F4CA N\u1ee3 hi\u1ec7n t\u1ea1i:                     0 VND",
            "\U0001F4B5 T\xedn d\u1ee5ng kh\u1ea3 d\u1ee5ng:         50,000,000 VND",
            "",
            "\U0001F4CB Thay \u0111\u1ed5i:",
            "\u2705 Ph\xe2n lo\u1ea1i: new \u2192 old",
            "\u2705 H\u1ea1n m\u1ee9c: 0 \u2192 50,000,000 VND",
            "",
            "KH ABC Tech gi\u1edd c\xf3 th\u1ec3 mua ch\u1ecbu t\u1ed1i \u0111a 50 tri\u1ec7u VND.",
        ], bot_height=Inches(3.8))

    # Slide 7: TC7
    make_test_slide(prs, 7, "Li\u1ec7t k\xea s\u1ea3n ph\u1ea9m",
        "Li\u1ec7t k\xea s\u1ea3n ph\u1ea9m c\xf3 trong h\u1ec7 th\u1ed1ng",
        [
            "\U0001F4E6 S\u1ea3n ph\u1ea9m trong h\u1ec7 th\u1ed1ng",
            "",
            "\u26A0\uFE0F Kh\xf4ng t\xecm th\u1ea5y s\u1ea3n ph\u1ea9m n\xe0o.",
            "",
            "H\u1ec7 th\u1ed1ng Odoo hi\u1ec7n t\u1ea1i ch\u01b0a c\xf3 s\u1ea3n ph\u1ea9m.",
            "",
            "\u0110\u1ec3 th\xeam s\u1ea3n ph\u1ea9m m\u1edbi:",
            "1. Th\xeam tr\u1ef1c ti\u1ebfp trong Odoo UI (Inventory \u2192 Products)",
            "2. Ho\u1eb7c t\xf4i c\xf3 th\u1ec3 t\u1ea1o s\u1ea3n ph\u1ea9m qua API",
            "",
            "B\u1ea1n mu\u1ed1n t\u1ea1o s\u1ea3n ph\u1ea9m m\u1edbi kh\xf4ng?",
        ], bot_height=Inches(3.2))

    # Slide 8: TC8
    make_test_slide(prs, 8, "T\u1ed5ng quan kinh doanh (Bot 1)",
        "T\u1ed5ng quan kinh doanh h\xf4m nay: KPI, pipeline, c\xf4ng n\u1ee3",
        [
            "\U0001F4CA T\u1ed5ng quan kinh doanh - 11/03/2026",
            "",
            "\U0001F3AF KPIs Th\xe1ng 3:",
            "\u2022 Leads m\u1edbi: 3",
            "\u2022 Gi\xe1 tr\u1ecb Pipeline: 0",
            "\u2022 Doanh thu th\xe1ng: 0",
            "\u2022 T\u1ed5ng c\xf4ng n\u1ee3: 0",
            "",
            "\U0001F4C8 Pipeline CRM:",
            "\U0001F195 New: 3  |  \u2705 Qualified: 0  |  \U0001F4CB Proposition: 0  |  \U0001F3C6 Won: 0",
            "",
            "\U0001F4B0 C\xf4ng n\u1ee3: \u2705 Kh\xf4ng c\xf3 KH v\u01b0\u1ee3t h\u1ea1n m\u1ee9c",
            "",
            "\U0001F4DD 3 leads ti\u1ec1m n\u0103ng c\u1ea7n qualify, kh\xf4ng c\xf3 c\xf4ng n\u1ee3 x\u1ea5u.",
        ], bot_height=Inches(3.8))

    # Slide 9: TC9
    make_test_slide(prs, 9, "T\xecm kh\xe1ch h\xe0ng",
        "T\xecm t\u1ea5t c\u1ea3 kh\xe1ch h\xe0ng, k\xe8m th\xf4ng tin c\xf4ng n\u1ee3",
        [
            "\U0001F465 Danh s\xe1ch Kh\xe1ch h\xe0ng",
            "",
            "\U0001F7E2 C\xf4ng ty ABC Tech (ID: 14) \u2014 Old \u2014 HM: 50,000,000 \u2014 N\u1ee3: 0",
            "\U0001F535 Trinity Masios (ID: 8) \u2014 New \u2014 HM: 0 \u2014 N\u1ee3: 0",
            "\U0001F535 Test KH MCP (ID: 10) \u2014 New \u2014 HM: 0 \u2014 N\u1ee3: 0",
            "\U0001F535 tran tan thanh (ID: 7) \u2014 New \u2014 HM: 0 \u2014 N\u1ee3: 0",
            "\U0001F535 Nguyen Van A (ID: 9) \u2014 New \u2014 HM: 0 \u2014 N\u1ee3: 0",
            "",
            "\U0001F4CA T\xf3m t\u1eaft:",
            "\U0001F7E2 Old: 1 KH \u2014 T\u1ed5ng HM: 50,000,000",
            "\U0001F535 New: 4+ KH \u2014 T\u1ed5ng HM: 0",
            "",
            "ABC Tech l\xe0 KH duy nh\u1ea5t c\xf3 h\u1ea1n m\u1ee9c c\xf4ng n\u1ee3.",
        ], bot_height=Inches(3.8))

    # Slide 10: TC10
    make_test_slide(prs, 10, "Kh\xe1ch h\xe0ng v\u01b0\u1ee3t h\u1ea1n m\u1ee9c",
        "C\xf3 kh\xe1ch h\xe0ng n\xe0o \u0111ang v\u01b0\u1ee3t h\u1ea1n m\u1ee9c c\xf4ng n\u1ee3 kh\xf4ng?",
        [
            "\u2705 Kh\xf4ng c\xf3 kh\xe1ch h\xe0ng n\xe0o v\u01b0\u1ee3t h\u1ea1n m\u1ee9c c\xf4ng n\u1ee3",
            "",
            "K\u1ebft qu\u1ea3: 0 kh\xe1ch h\xe0ng \u0111ang v\u01b0\u1ee3t h\u1ea1n m\u1ee9c.",
            "",
            "C\xf4ng ty ABC Tech \u2014 HM: 50,000,000 \u2014 N\u1ee3: 0 \u2014 \u2705 B\xecnh th\u01b0\u1eddng",
            "",
            "T\u1ea5t c\u1ea3 kh\xe1ch h\xe0ng \u0111\u1ec1u trong t\xecnh tr\u1ea1ng t\xe0i ch\xednh t\u1ed1t.",
        ], bot_height=Inches(2.4))

    # Slide 11: Summary
    make_summary(prs)

    out_path = os.path.join(os.path.dirname(__file__), "testcase-slides.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
